#!/usr/bin/env python3
"""
SharePoint Integration Module for KnowledgeOps Agent
Enterprise-grade integration with Microsoft Graph API, document processing, and managed metadata support.
"""

import os
import json
import logging
import asyncio
import aiohttp
import time
import io
from datetime import datetime, timezone
from typing import Dict, List, Optional, Any, Union, BinaryIO
from dataclasses import dataclass, asdict
from urllib.parse import urljoin, quote, urlparse
import base64
import mimetypes
from pathlib import Path
import zipfile
import tempfile

# Document processing libraries
from docx import Document
import openpyxl
from pptx import Presentation
import PyPDF2
import xml.etree.ElementTree as ET

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class SharePointConfig:
    """Configuration for SharePoint integration"""
    tenant_id: str
    client_id: str
    client_secret: str
    auth_type: str = 'client_credentials'  # 'client_credentials', 'authorization_code'
    graph_endpoint: str = 'https://graph.microsoft.com/v1.0'
    sites: Optional[List[str]] = None
    document_libraries: Optional[List[str]] = None
    max_results: int = 100
    timeout: int = 30
    retry_attempts: int = 3
    rate_limit_per_minute: int = 2000
    supported_file_types: List[str] = None

    def __post_init__(self):
        if self.supported_file_types is None:
            self.supported_file_types = [
                'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 
                'pdf', 'txt', 'html', 'htm', 'md'
            ]

@dataclass
class SharePointContent:
    """Standardized content structure from SharePoint"""
    id: str
    name: str
    title: str
    content_text: str
    summary: str
    author: str
    created_date: datetime
    modified_date: datetime
    site_id: str
    site_name: str
    library_name: str
    url: str
    web_url: str
    content_type: str
    file_type: str
    file_size: int
    tags: List[str]
    managed_metadata: Dict[str, Any]
    permissions: Dict[str, List[str]]
    metadata: Dict[str, Any]
    parent_folder: Optional[str] = None
    version: str = "1.0"

class SharePointAuthManager:
    """Handles authentication with Microsoft Graph API"""
    
    def __init__(self, config: SharePointConfig):
        self.config = config
        self.session = aiohttp.ClientSession(
            timeout=aiohttp.ClientTimeout(total=config.timeout)
        )
        self._access_token = None
        self._token_expires_at = None
        
    async def authenticate(self) -> Dict[str, str]:
        """Authenticate and get access token"""
        if self._access_token and self._token_expires_at > time.time():
            return {'Authorization': f'Bearer {self._access_token}'}
        
        if self.config.auth_type == 'client_credentials':
            return await self._client_credentials_auth()
        else:
            raise ValueError(f"Unsupported auth type: {self.config.auth_type}")
    
    async def _client_credentials_auth(self) -> Dict[str, str]:
        """Client credentials flow for service-to-service authentication"""
        token_url = f'https://login.microsoftonline.com/{self.config.tenant_id}/oauth2/v2.0/token'
        
        data = {
            'client_id': self.config.client_id,
            'client_secret': self.config.client_secret,
            'scope': 'https://graph.microsoft.com/.default',
            'grant_type': 'client_credentials'
        }
        
        async with self.session.post(token_url, data=data) as response:
            if response.status == 200:
                token_data = await response.json()
                self._access_token = token_data['access_token']
                self._token_expires_at = time.time() + token_data.get('expires_in', 3600) - 60
                
                logger.info("Microsoft Graph authentication successful")
                return {'Authorization': f'Bearer {self._access_token}'}
            else:
                error_text = await response.text()
                raise Exception(f"Authentication failed: {response.status} - {error_text}")
    
    async def close(self):
        """Close the session"""
        await self.session.close()

class SharePointDocumentProcessor:
    """Processes different document types from SharePoint"""
    
    def __init__(self):
        self.processors = {
            'docx': self._process_docx,
            'doc': self._process_doc,
            'xlsx': self._process_xlsx,
            'xls': self._process_xls,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'pdf': self._process_pdf,
            'txt': self._process_text,
            'html': self._process_html,
            'htm': self._process_html,
            'md': self._process_markdown
        }
    
    async def process_document(self, file_content: bytes, file_extension: str, filename: str) -> Dict[str, Any]:
        """Process document content based on file type"""
        processor = self.processors.get(file_extension.lower())
        
        if not processor:
            logger.warning(f"No processor available for file type: {file_extension}")
            return {
                'text_content': '',
                'metadata': {'error': f'Unsupported file type: {file_extension}'}
            }
        
        try:
            return await processor(file_content, filename)
        except Exception as e:
            logger.error(f"Error processing {filename}: {e}")
            return {
                'text_content': '',
                'metadata': {'error': str(e)}
            }
    
    async def _process_docx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Word DOCX files"""
        with tempfile.NamedTemporaryFile() as temp_file:
            temp_file.write(content)
            temp_file.flush()
            
            doc = Document(temp_file.name)
            
            # Extract text
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract tables
            table_content = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_content.append(' | '.join(row_text))
            
            # Combine content
            full_text = '\n'.join(text_content)
            if table_content:
                full_text += '\n\nTables:\n' + '\n'.join(table_content)
            
            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                'author': core_props.author or '',
                'title': core_props.title or '',
                'subject': core_props.subject or '',
                'created': core_props.created.isoformat() if core_props.created else '',
                'modified': core_props.modified.isoformat() if core_props.modified else '',
                'word_count': len(full_text.split()),
                'paragraph_count': len(text_content),
                'table_count': len(table_content)
            }
            
            return {
                'text_content': full_text,
                'metadata': metadata
            }
    
    async def _process_doc(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process legacy Word DOC files"""
        # For legacy DOC files, we'd typically use python-docx2txt or similar
        # This is a simplified implementation
        return {
            'text_content': '[Legacy DOC file - content extraction requires additional libraries]',
            'metadata': {'file_type': 'legacy_doc'}
        }
    
    async def _process_xlsx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Excel XLSX files"""
        with tempfile.NamedTemporaryFile() as temp_file:
            temp_file.write(content)
            temp_file.flush()
            
            workbook = openpyxl.load_workbook(temp_file.name, data_only=True)
            
            text_content = []
            sheet_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(cell.strip() for cell in row_data):
                        sheet_content.append(' | '.join(row_data))
                
                if sheet_content:
                    sheet_data[sheet_name] = sheet_content
                    text_content.extend([f"Sheet: {sheet_name}"] + sheet_content)
            
            metadata = {
                'sheet_count': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames,
                'total_rows': sum(len(data) for data in sheet_data.values())
            }
            
            return {
                'text_content': '\n'.join(text_content),
                'metadata': metadata
            }
    
    async def _process_xls(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process legacy Excel XLS files"""
        # For legacy XLS files, we'd typically use xlrd or similar
        return {
            'text_content': '[Legacy XLS file - content extraction requires additional libraries]',
            'metadata': {'file_type': 'legacy_xls'}
        }
    
    async def _process_pptx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint PPTX files"""
        with tempfile.NamedTemporaryFile() as temp_file:
            temp_file.write(content)
            temp_file.flush()
            
            presentation = Presentation(temp_file.name)
            
            text_content = []
            slide_count = 0
            
            for slide in presentation.slides:
                slide_count += 1
                slide_text = [f"Slide {slide_count}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide number
                    text_content.extend(slide_text)
            
            metadata = {
                'slide_count': slide_count,
                'total_shapes': sum(len(slide.shapes) for slide in presentation.slides)
            }
            
            return {
                'text_content': '\n'.join(text_content),
                'metadata': metadata
            }
    
    async def _process_ppt(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process legacy PowerPoint PPT files"""
        return {
            'text_content': '[Legacy PPT file - content extraction requires additional libraries]',
            'metadata': {'file_type': 'legacy_ppt'}
        }
    
    async def _process_pdf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF files"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(content)
                temp_file.flush()
                
                with open(temp_file.name, 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    
                    text_content = []
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text.strip():
                                text_content.append(f"Page {page_num + 1}:\n{page_text.strip()}")
                        except Exception as e:
                            logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    
                    metadata = {
                        'page_count': len(pdf_reader.pages),
                        'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                        'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                        'creator': pdf_reader.metadata.get('/Creator', '') if pdf_reader.metadata else ''
                    }
                    
                    return {
                        'text_content': '\n\n'.join(text_content),
                        'metadata': metadata
                    }
        except Exception as e:
            logger.error(f"Error processing PDF {filename}: {e}")
            return {
                'text_content': '[PDF processing failed]',
                'metadata': {'error': str(e)}
            }
    
    async def _process_text(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            text_content = content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                text_content = content.decode('latin-1')
            except UnicodeDecodeError:
                text_content = content.decode('utf-8', errors='ignore')
        
        metadata = {
            'character_count': len(text_content),
            'line_count': len(text_content.splitlines()),
            'word_count': len(text_content.split())
        }
        
        return {
            'text_content': text_content,
            'metadata': metadata
        }
    
    async def _process_html(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process HTML files"""
        from bs4 import BeautifulSoup
        
        try:
            html_content = content.decode('utf-8')
        except UnicodeDecodeError:
            html_content = content.decode('utf-8', errors='ignore')
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Extract text
        text_content = soup.get_text(separator=' ', strip=True)
        
        # Extract metadata
        title = soup.find('title')
        meta_description = soup.find('meta', attrs={'name': 'description'})
        meta_keywords = soup.find('meta', attrs={'name': 'keywords'})
        
        metadata = {
            'title': title.get_text() if title else '',
            'description': meta_description.get('content', '') if meta_description else '',
            'keywords': meta_keywords.get('content', '') if meta_keywords else '',
            'word_count': len(text_content.split())
        }
        
        return {
            'text_content': text_content,
            'metadata': metadata
        }
    
    async def _process_markdown(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Markdown files"""
        try:
            text_content = content.decode('utf-8')
        except UnicodeDecodeError:
            text_content = content.decode('utf-8', errors='ignore')
        
        # Basic markdown processing - extract headers and content
        lines = text_content.splitlines()
        headers = []
        
        for line in lines:
            if line.startswith('#'):
                headers.append(line.strip())
        
        metadata = {
            'header_count': len(headers),
            'headers': headers[:10],  # First 10 headers
            'line_count': len(lines),
            'word_count': len(text_content.split())
        }
        
        return {
            'text_content': text_content,
            'metadata': metadata
        }

class SharePointManagedMetadataManager:
    """Manages SharePoint managed metadata and taxonomy"""
    
    def __init__(self, auth_manager):
        self.auth_manager = auth_manager
        
    async def get_term_store_terms(self, term_set_id: str) -> List[Dict[str, Any]]:
        """Get terms from a specific term set"""
        headers = await self.auth_manager.authenticate()
        
        terms_url = f"{self.auth_manager.config.graph_endpoint}/termStore/sets/{term_set_id}/terms"
        
        async with self.auth_manager.session.get(terms_url, headers=headers) as response:
            if response.status == 200:
                data = await response.json()
                return data.get('value', [])
            else:
                logger.warning(f"Failed to get terms for set {term_set_id}: {response.status}")
                return []
    
    async def resolve_managed_metadata_field(self, field_value: Dict) -> Dict[str, Any]:
        """Resolve managed metadata field to human-readable format"""
        if not isinstance(field_value, dict) or 'TermGuid' not in field_value:
            return {'label': str(field_value), 'term_path': []}
        
        term_guid = field_value['TermGuid']
        headers = await self.auth_manager.authenticate()
        
        term_url = f"{self.auth_manager.config.graph_endpoint}/termStore/terms/{term_guid}"
        
        async with self.auth_manager.session.get(term_url, headers=headers) as response:
            if response.status == 200:
                term_data = await response.json()
                return {
                    'label': field_value.get('Label', ''),
                    'term_id': term_guid,
                    'term_path': await self._build_term_path(term_data),
                    'description': term_data.get('descriptions', [{}])[0].get('description', '')
                }
            else:
                return {
                    'label': field_value.get('Label', ''),
                    'term_id': term_guid,
                    'term_path': [],
                    'error': f"Failed to resolve term: {response.status}"
                }
    
    async def _build_term_path(self, term_data: Dict) -> List[str]:
        """Build hierarchical path for a term"""
        path = [term_data.get('labels', [{}])[0].get('name', '')]
        
        # Get parent terms if available
        parent_id = term_data.get('parent', {}).get('id')
        if parent_id:
            headers = await self.auth_manager.authenticate()
            parent_url = f"{self.auth_manager.config.graph_endpoint}/termStore/terms/{parent_id}"
            
            async with self.auth_manager.session.get(parent_url, headers=headers) as response:
                if response.status == 200:
                    parent_data = await response.json()
                    parent_path = await self._build_term_path(parent_data)
                    path = parent_path + path
        
        return path

class SharePointAPIClient:
    """Main SharePoint API client with enterprise features"""
    
    def __init__(self, config: SharePointConfig):
        self.config = config
        self.auth_manager = SharePointAuthManager(config)
        self.document_processor = SharePointDocumentProcessor()
        self.metadata_manager = SharePointManagedMetadataManager(self.auth_manager)
        self._rate_limiter = asyncio.Semaphore(config.rate_limit_per_minute // 60)
        
    async def initialize(self):
        """Initialize the client and authenticate"""
        await self.auth_manager.authenticate()
        logger.info("SharePoint API client initialized")
    
    async def discover_sites(self) -> List[Dict[str, Any]]:
        """Discover all SharePoint sites in the tenant"""
        headers = await self.auth_manager.authenticate()
        
        sites_url = f"{self.config.graph_endpoint}/sites"
        params = {
            '$select': 'id,name,webUrl,description,createdDateTime,lastModifiedDateTime',
            '$top': self.config.max_results
        }
        
        all_sites = []
        
        async with self._rate_limiter:
            async with self.auth_manager.session.get(
                sites_url, headers=headers, params=params
            ) as response:
                if response.status == 200:
                    data = await response.json()
                    sites = data.get('value', [])
                    all_sites.extend(sites)
                    
                    # Handle pagination
                    next_link = data.get('@odata.nextLink')
                    while next_link:
                        async with self.auth_manager.session.get(
                            next_link, headers=headers
                        ) as next_response:
                            if next_response.status == 200:
                                next_data = await next_response.json()
                                sites = next_data.get('value', [])
                                all_sites.extend(sites)
                                next_link = next_data.get('@odata.nextLink')
                            else:
                                break
                else:
                    error_text = await response.text()
                    logger.error(f"Failed to discover sites: {response.status} - {error_text}")
        
        logger.info(f"Discovered {len(all_sites)} sites")
        return all_sites
    
    async def get_site_document_libraries(self, site_id: str) -> List[Dict[str, Any]]:
        """Get document libraries for a specific site"""
        headers = await self.auth_manager.authenticate()
        
        drives_url = f"{self.config.graph_endpoint}/sites/{site_id}/drives"
        params = {
            '$select': 'id,name,description,driveType,quota',
            '$filter': "driveType eq 'documentLibrary'"
        }
        
        async with self._rate_limiter:
            async with self.auth_manager.session.get(
                drives_url, headers=headers, params=params
            ) as response:
                if response.status == 200:
                    data = await response.json()
                    return data.get('value', [])
                else:
                    logger.warning(f"Failed to get document libraries for site {site_id}: {response.status}")
                    return []
    
    async def search_content(self, query: str, site_ids: Optional[List[str]] = None) -> List[SharePointContent]:
        """Search content across SharePoint sites"""
        headers = await self.auth_manager.authenticate()
        headers['Content-Type'] = 'application/json'
        
        search_url = f"{self.config.graph_endpoint}/search/query"
        
        # Build search request
        search_request = {
            'requests': [{
                'entityTypes': ['driveItem'],
                'query': {
                    'queryString': query
                },
                'from': 0,
                'size': self.config.max_results,
                'fields': [
                    'id', 'name', 'webUrl', 'lastModifiedDateTime', 
                    'createdDateTime', 'size', 'fileSystemInfo',
                    'createdBy', 'lastModifiedBy', 'parentReference'
                ]
            }]
        }
        
        # Add site filter if specified
        if site_ids:
            site_filter = ' OR '.join([f'SiteId:{site_id}' for site_id in site_ids])
            search_request['requests'][0]['query']['queryString'] += f' AND ({site_filter})'
        
        all_results = []
        
        async with self._rate_limiter:
            async with self.auth_manager.session.post(
                search_url, headers=headers, json=search_request
            ) as response:
                if response.status == 200:
                    data = await response.json()
                    search_results = data.get('value', [{}])[0].get('hitsContainers', [{}])[0].get('hits', [])
                    
                    for hit in search_results:
                        try:
                            content = await self._process_search_result(hit)
                            if content:
                                all_results.append(content)
                        except Exception as e:
                            logger.error(f"Error processing search result: {e}")
                else:
                    error_text = await response.text()
                    logger.error(f"Search failed: {response.status} - {error_text}")
        
        logger.info(f"Found {len(all_results)} content items")
        return all_results
    
    async def get_document_content(self, site_id: str, drive_id: str, item_id: str) -> Optional[SharePointContent]:
        """Get detailed content for a specific document"""
        headers = await self.auth_manager.authenticate()
        
        # Get item metadata
        item_url = f"{self.config.graph_endpoint}/sites/{site_id}/drives/{drive_id}/items/{item_id}"
        params = {'expand': 'listItem($expand=fields)'}
        
        async with self._rate_limiter:
            async with self.auth_manager.session.get(
                item_url, headers=headers, params=params
            ) as response:
                if response.status == 200:
                    item_data = await response.json()
                    return await self._process_document_item(item_data, site_id)
                else:
                    logger.error(f"Failed to get document {item_id}: {response.status}")
                    return None
    
    async def sync_site_content(self, site_id: str, since: Optional[datetime] = None) -> List[SharePointContent]:
        """Sync all content from a specific site"""
        headers = await self.auth_manager.authenticate()
        
        # Get document libraries
        libraries = await self.get_site_document_libraries(site_id)
        
        all_content = []
        
        for library in libraries:
            drive_id = library['id']
            library_name = library['name']
            
            logger.info(f"Syncing content from library: {library_name}")
            
            # Get all items in the library
            items_url = f"{self.config.graph_endpoint}/sites/{site_id}/drives/{drive_id}/root/children"
            params = {
                '$expand': 'listItem($expand=fields)',
                '$top': self.config.max_results
            }
            
            if since:
                since_str = since.isoformat()
                params['$filter'] = f"lastModifiedDateTime ge {since_str}"
            
            try:
                library_content = await self._get_library_items_recursive(
                    site_id, drive_id, items_url, params, headers
                )
                all_content.extend(library_content)
                logger.info(f"Synced {len(library_content)} items from library {library_name}")
            except Exception as e:
                logger.error(f"Error syncing library {library_name}: {e}")
        
        logger.info(f"Total content synced from site: {len(all_content)} items")
        return all_content
    
    async def _get_library_items_recursive(self, site_id: str, drive_id: str, url: str, params: Dict, headers: Dict) -> List[SharePointContent]:
        """Recursively get all items from a library including subfolders"""
        items = []
        
        async with self._rate_limiter:
            async with self.auth_manager.session.get(url, headers=headers, params=params) as response:
                if response.status == 200:
                    data = await response.json()
                    
                    for item in data.get('value', []):
                        try:
                            if item.get('file'):  # It's a file
                                content = await self._process_document_item(item, site_id)
                                if content:
                                    items.append(content)
                            elif item.get('folder'):  # It's a folder
                                # Recursively process folder contents
                                folder_url = f"{self.config.graph_endpoint}/sites/{site_id}/drives/{drive_id}/items/{item['id']}/children"
                                folder_items = await self._get_library_items_recursive(
                                    site_id, drive_id, folder_url, params, headers
                                )
                                items.extend(folder_items)
                        except Exception as e:
                            logger.error(f"Error processing item {item.get('id', 'unknown')}: {e}")
                    
                    # Handle pagination
                    next_link = data.get('@odata.nextLink')
                    if next_link:
                        next_items = await self._get_library_items_recursive(
                            site_id, drive_id, next_link, {}, headers
                        )
                        items.extend(next_items)
                else:
                    logger.warning(f"Failed to get library items: {response.status}")
        
        return items
    
    async def _process_search_result(self, hit: Dict) -> Optional[SharePointContent]:
        """Process a search result hit into SharePointContent"""
        try:
            resource = hit.get('resource', {})
            
            # Extract basic information
            item_id = resource.get('id', '')
            name = resource.get('name', '')
            web_url = resource.get('webUrl', '')
            
            # Parse site and drive info from web URL
            parsed_url = urlparse(web_url)
            site_info = await self._extract_site_info_from_url(web_url)
            
            # Create basic content object
            content = SharePointContent(
                id=item_id,
                name=name,
                title=name,
                content_text='',  # Will be filled by document processing
                summary='',
                author=resource.get('createdBy', {}).get('user', {}).get('displayName', ''),
                created_date=datetime.fromisoformat(resource.get('createdDateTime', '').replace('Z', '+00:00')),
                modified_date=datetime.fromisoformat(resource.get('lastModifiedDateTime', '').replace('Z', '+00:00')),
                site_id=site_info.get('site_id', ''),
                site_name=site_info.get('site_name', ''),
                library_name=site_info.get('library_name', ''),
                url=web_url,
                web_url=web_url,
                content_type='document',
                file_type=Path(name).suffix.lstrip('.').lower(),
                file_size=resource.get('size', 0),
                tags=[],
                managed_metadata={},
                permissions={},
                metadata={}
            )
            
            return content
            
        except Exception as e:
            logger.error(f"Error processing search result: {e}")
            return None
    
    async def _process_document_item(self, item_data: Dict, site_id: str) -> Optional[SharePointContent]:
        """Process a document item into SharePointContent with full content extraction"""
        try:
            # Extract basic information
            item_id = item_data['id']
            name = item_data['name']
            web_url = item_data.get('webUrl', '')
            file_info = item_data.get('file', {})
            
            # Get file extension
            file_extension = Path(name).suffix.lstrip('.').lower()
            
            # Skip unsupported file types
            if file_extension not in self.config.supported_file_types:
                logger.debug(f"Skipping unsupported file type: {file_extension}")
                return None
            
            # Download and process file content
            content_text = ''
            document_metadata = {}
            
            if file_info and '@microsoft.graph.downloadUrl' in item_data:
                download_url = item_data['@microsoft.graph.downloadUrl']
                
                async with self._rate_limiter:
                    async with self.auth_manager.session.get(download_url) as response:
                        if response.status == 200:
                            file_content = await response.read()
                            
                            # Process document content
                            processed_content = await self.document_processor.process_document(
                                file_content, file_extension, name
                            )
                            content_text = processed_content.get('text_content', '')
                            document_metadata = processed_content.get('metadata', {})
            
            # Generate summary
            summary = content_text[:300] + '...' if len(content_text) > 300 else content_text
            
            # Extract SharePoint-specific metadata
            list_item = item_data.get('listItem', {})
            fields = list_item.get('fields', {})
            
            # Process managed metadata fields
            managed_metadata = {}
            for field_name, field_value in fields.items():
                if isinstance(field_value, dict) and 'TermGuid' in field_value:
                    managed_metadata[field_name] = await self.metadata_manager.resolve_managed_metadata_field(field_value)
            
            # Extract tags from fields
            tags = []
            if 'Tags' in fields:
                tags_field = fields['Tags']
                if isinstance(tags_field, list):
                    tags = [tag.get('Label', '') for tag in tags_field if isinstance(tag, dict)]
                elif isinstance(tags_field, str):
                    tags = [tag.strip() for tag in tags_field.split(';') if tag.strip()]
            
            # Get site information
            site_info = await self._get_site_info(site_id)
            
            # Create content object
            content = SharePointContent(
                id=item_id,
                name=name,
                title=fields.get('Title', name),
                content_text=content_text,
                summary=summary,
                author=item_data.get('createdBy', {}).get('user', {}).get('displayName', ''),
                created_date=datetime.fromisoformat(item_data.get('createdDateTime', '').replace('Z', '+00:00')),
                modified_date=datetime.fromisoformat(item_data.get('lastModifiedDateTime', '').replace('Z', '+00:00')),
                site_id=site_id,
                site_name=site_info.get('name', ''),
                library_name=item_data.get('parentReference', {}).get('name', ''),
                url=web_url,
                web_url=web_url,
                content_type='document',
                file_type=file_extension,
                file_size=item_data.get('size', 0),
                tags=tags,
                managed_metadata=managed_metadata,
                permissions={},  # Would need additional API calls to get detailed permissions
                metadata={
                    **document_metadata,
                    'sharepoint_fields': fields,
                    'mime_type': file_info.get('mimeType', ''),
                    'hash': file_info.get('hashes', {})
                },
                parent_folder=item_data.get('parentReference', {}).get('path', ''),
                version=str(item_data.get('eTag', '1.0'))
            )
            
            return content
            
        except Exception as e:
            logger.error(f"Error processing document item {item_data.get('id', 'unknown')}: {e}")
            return None
    
    async def _get_site_info(self, site_id: str) -> Dict[str, Any]:
        """Get basic site information"""
        headers = await self.auth_manager.authenticate()
        
        site_url = f"{self.config.graph_endpoint}/sites/{site_id}"
        
        async with self._rate_limiter:
            async with self.auth_manager.session.get(site_url, headers=headers) as response:
                if response.status == 200:
                    return await response.json()
                else:
                    return {'name': 'Unknown Site'}
    
    async def _extract_site_info_from_url(self, web_url: str) -> Dict[str, str]:
        """Extract site and library information from SharePoint URL"""
        # This is a simplified implementation
        # In practice, you'd parse the URL structure to extract site and library info
        return {
            'site_id': 'unknown',
            'site_name': 'Unknown Site',
            'library_name': 'Unknown Library'
        }
    
    async def close(self):
        """Close the client and cleanup resources"""
        await self.auth_manager.close()
        logger.info("SharePoint API client closed")

# Example usage and configuration
async def main():
    """Example usage of the SharePoint integration"""
    
    # Configuration for SharePoint Online
    config = SharePointConfig(
        tenant_id=os.getenv("SHAREPOINT_TENANT_ID"),
        client_id=os.getenv("SHAREPOINT_CLIENT_ID"),
        client_secret=os.getenv("SHAREPOINT_CLIENT_SECRET"),
        sites=["https://company.sharepoint.com/sites/engineering"],
        max_results=50
    )
    
    client = SharePointAPIClient(config)
    
    try:
        await client.initialize()
        
        # Discover sites
        sites = await client.discover_sites()
        print(f"Found {len(sites)} sites")
        
        # Search for content
        search_results = await client.search_content("deployment procedures")
        print(f"Found {len(search_results)} documents about deployment procedures")
        
        # Sync content from a specific site
        if sites:
            site_id = sites[0]['id']
            site_content = await client.sync_site_content(site_id)
            print(f"Synced {len(site_content)} documents from site")
            
            # Example: Get specific document content
            if site_content:
                doc = site_content[0]
                print(f"Document: {doc.title}")
                print(f"Content length: {len(doc.content_text)} characters")
                print(f"File type: {doc.file_type}")
                print(f"Tags: {doc.tags}")
                print(f"Managed metadata: {doc.managed_metadata}")
        
    finally:
        await client.close()

if __name__ == "__main__":
    asyncio.run(main())

